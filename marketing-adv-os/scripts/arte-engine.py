#!/usr/bin/env python3
"""
arte-engine.py — Engine reutilizavel de geracao de artes (PNG + PPTX).

Le tokens.json (paleta + tipografia) + sizes.json (mapa de medidas) + template.json
(layout) + dados do operador (texto, paths de imagens), e produz:
- PNG-24 sRGB com texto renderizado nas regioes definidas no template
- PPTX equivalente (1 slide) pronto para abrir no Canva ou editor de slides

Verifica disponibilidade de Pillow e python-pptx ANTES de executar. Se faltam:
- Pillow ausente -> gera arquivo SVG de fallback (universal, sem deps)
- python-pptx ausente -> gera HTML standalone (que pode ser convertido em PPT via Canva/LibreOffice)

Uso direto:
    python3 arte-engine.py --template feed-portrait --texto "Headline aqui" \\
        --subtexto "Subhead opcional" --output ./artes/peca-01.png \\
        --tokens ./marketing/design-tokens/tokens.json \\
        --sizes ./templates/arte-data/sizes.json

Uso como modulo:
    from arte_engine import render_png, render_pptx, check_deps
    deps = check_deps()
    if deps['pillow']:
        render_png(template, tokens, texto, output_path)

Stdlib + Pillow (opcional) + python-pptx (opcional).
"""

from __future__ import annotations

import argparse
import io
import json
import sys
from pathlib import Path
from typing import Optional


# ============================================================
# Verificacao de deps
# ============================================================

def check_deps() -> dict:
    """Retorna dict {'pillow': bool, 'python-pptx': bool, 'requests': bool}."""
    deps = {}
    try:
        import PIL  # noqa: F401
        deps['pillow'] = True
    except ImportError:
        deps['pillow'] = False
    try:
        import pptx  # noqa: F401
        deps['python-pptx'] = True
    except ImportError:
        deps['python-pptx'] = False
    try:
        import requests  # noqa: F401
        deps['requests'] = True
    except ImportError:
        deps['requests'] = False
    return deps


def print_deps_status(deps: dict) -> None:
    print("Dependencias detectadas:", file=sys.stderr)
    for name, ok in deps.items():
        marker = "OK" if ok else "FALTANDO"
        print(f"  [{marker}] {name}", file=sys.stderr)
    if not all(deps.values()):
        print("\nPara instalar: pip3 install pillow python-pptx requests", file=sys.stderr)


# ============================================================
# Carregar tokens + sizes + template
# ============================================================

def load_tokens(tokens_path: Path) -> dict:
    """Le tokens.json gerado por marketing-tokens-css."""
    if not tokens_path.exists():
        raise FileNotFoundError(
            f"tokens.json nao encontrado em {tokens_path}. "
            f"Rode `marketing-tokens-css` primeiro para gerar."
        )
    return json.loads(tokens_path.read_text(encoding="utf-8"))


def load_sizes(sizes_path: Path) -> dict:
    """Le sizes.json do plugin (mapa canonico de medidas)."""
    if not sizes_path.exists():
        raise FileNotFoundError(f"sizes.json nao encontrado em {sizes_path}")
    return json.loads(sizes_path.read_text(encoding="utf-8"))


def get_canvas_spec(sizes: dict, template_key: str) -> dict:
    """
    Resolve template_key (ex: 'instagram.feed_portrait') -> spec do canvas.
    Retorna dict com canvas + ai_safe_v3 (ou v2 conforme disponivel).
    """
    parts = template_key.split('.')
    if len(parts) != 2:
        raise ValueError(f"template_key deve ser 'plataforma.formato', recebi {template_key!r}")
    plataforma, formato = parts
    if plataforma not in sizes:
        raise KeyError(f"Plataforma '{plataforma}' nao encontrada em sizes.json")
    if formato not in sizes[plataforma]:
        raise KeyError(f"Formato '{formato}' nao encontrado em sizes.{plataforma}")
    return sizes[plataforma][formato]


# ============================================================
# Render PNG (Pillow path)
# ============================================================

def render_png(
    canvas_spec: dict,
    tokens: dict,
    texto: str,
    subtexto: Optional[str],
    output_path: Path,
    layout: str = "centered",
) -> Path:
    """
    Renderiza PNG usando Pillow. Layouts: centered, top-aligned, hero-with-subtext.
    Aplica paleta + tipografia do tokens.json. Respeita AI-SAFE zone.

    Returns:
        Path para o PNG gerado.
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        raise RuntimeError("Pillow nao instalado. `pip3 install pillow`")

    canvas = canvas_spec['canvas']
    width, height = canvas['width'], canvas['height']
    palette = tokens.get('palette', {})
    bg_color = palette.get('neutral_light', '#FFFFFF')
    fg_color = palette.get('neutral_dark', '#000000')
    accent_color = palette.get('accent', fg_color)

    # Criar canvas (RGB sRGB por default no Pillow)
    img = Image.new('RGB', (width, height), bg_color)
    draw = ImageDraw.Draw(img)

    # Determinar zona segura (V3 prioritario, fallback V2)
    safe_zone = canvas_spec.get('ai_safe_v3') or canvas_spec.get('ai_safe_v2')
    if safe_zone:
        sz_w = safe_zone['width']
        sz_h = safe_zone['height']
        sz_x = (width - sz_w) // 2
        sz_y = safe_zone.get('margin_top', (height - sz_h) // 2)
    else:
        # Fallback: margem padrao 10%
        sz_w = int(width * 0.8)
        sz_h = int(height * 0.8)
        sz_x = int(width * 0.1)
        sz_y = int(height * 0.1)

    # Carregar fonte — usar primary do tokens se possivel; fallback sistema
    try:
        scale = tokens.get('scale', {})
        h1_size = scale.get('h1', {}).get('size', 64)
        body_size = scale.get('body', {}).get('size', 28)
        # Tentar fonte primaria (precisa estar disponivel no sistema OU pre-baixada)
        try:
            font_primary = ImageFont.truetype("Helvetica", h1_size)
            font_secondary = ImageFont.truetype("Helvetica", body_size)
        except OSError:
            font_primary = ImageFont.load_default()
            font_secondary = ImageFont.load_default()
    except Exception:
        font_primary = ImageFont.load_default()
        font_secondary = ImageFont.load_default()

    # Layout centered: titulo grande no centro vertical da AI-SAFE
    if layout == "centered":
        # Quebrar texto em multiplas linhas se necessario
        lines = _wrap_text(draw, texto, font_primary, sz_w)
        line_height = h1_size + 20
        total_h = line_height * len(lines)
        y = sz_y + (sz_h - total_h) // 2
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=font_primary)
            text_w = bbox[2] - bbox[0]
            x = sz_x + (sz_w - text_w) // 2
            draw.text((x, y), line, fill=fg_color, font=font_primary)
            y += line_height

        # Subtexto opcional logo abaixo
        if subtexto:
            sub_lines = _wrap_text(draw, subtexto, font_secondary, sz_w)
            for line in sub_lines:
                bbox = draw.textbbox((0, 0), line, font=font_secondary)
                text_w = bbox[2] - bbox[0]
                x = sz_x + (sz_w - text_w) // 2
                draw.text((x, y + 20), line, fill=accent_color, font=font_secondary)
                y += body_size + 10

    elif layout == "top-aligned":
        # Texto no topo da AI-SAFE
        y = sz_y + 40
        lines = _wrap_text(draw, texto, font_primary, sz_w)
        for line in lines:
            draw.text((sz_x, y), line, fill=fg_color, font=font_primary)
            y += h1_size + 20

    elif layout == "hero-with-subtext":
        # Texto principal no terco superior, subtexto no terco inferior
        title_lines = _wrap_text(draw, texto, font_primary, sz_w)
        title_y = sz_y + sz_h // 4
        for line in title_lines:
            bbox = draw.textbbox((0, 0), line, font=font_primary)
            text_w = bbox[2] - bbox[0]
            x = sz_x + (sz_w - text_w) // 2
            draw.text((x, title_y), line, fill=fg_color, font=font_primary)
            title_y += h1_size + 20
        if subtexto:
            sub_lines = _wrap_text(draw, subtexto, font_secondary, sz_w)
            sub_y = sz_y + (sz_h * 3) // 4
            for line in sub_lines:
                bbox = draw.textbbox((0, 0), line, font=font_secondary)
                text_w = bbox[2] - bbox[0]
                x = sz_x + (sz_w - text_w) // 2
                draw.text((x, sub_y), line, fill=accent_color, font=font_secondary)
                sub_y += body_size + 10

    output_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(output_path, "PNG", optimize=True)
    return output_path


def _wrap_text(draw, text: str, font, max_width: int) -> list[str]:
    """Quebra texto em linhas que cabem em max_width."""
    words = text.split()
    if not words:
        return [text]
    lines = []
    current = words[0]
    for w in words[1:]:
        test = current + " " + w
        bbox = draw.textbbox((0, 0), test, font=font)
        if (bbox[2] - bbox[0]) <= max_width:
            current = test
        else:
            lines.append(current)
            current = w
    lines.append(current)
    return lines


# ============================================================
# Render PPTX (python-pptx path)
# ============================================================

def render_pptx(
    canvas_spec: dict,
    tokens: dict,
    texto: str,
    subtexto: Optional[str],
    output_path: Path,
) -> Path:
    """Renderiza PPTX equivalente para abrir em Canva/LibreOffice/PowerPoint."""
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise RuntimeError("python-pptx nao instalado. `pip3 install python-pptx`")

    canvas = canvas_spec['canvas']
    # PPTX usa English Metric Units (EMU). 1 px = 9525 EMU em 96 dpi.
    width_emu = Emu(canvas['width'] * 9525)
    height_emu = Emu(canvas['height'] * 9525)

    prs = Presentation()
    prs.slide_width = width_emu
    prs.slide_height = height_emu

    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    palette = tokens.get('palette', {})
    bg_hex = palette.get('neutral_light', '#FFFFFF').lstrip('#')
    fg_hex = palette.get('neutral_dark', '#000000').lstrip('#')

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(bg_hex)

    # Text box no centro da AI-SAFE zone
    safe = canvas_spec.get('ai_safe_v3') or canvas_spec.get('ai_safe_v2')
    if safe:
        tx_w = Emu(safe['width'] * 9525)
        tx_h = Emu(safe['height'] * 9525)
        tx_x = Emu((canvas['width'] - safe['width']) // 2 * 9525)
        tx_y = Emu(safe.get('margin_top', (canvas['height'] - safe['height']) // 2) * 9525)
    else:
        tx_w = Emu(int(canvas['width'] * 0.8) * 9525)
        tx_h = Emu(int(canvas['height'] * 0.8) * 9525)
        tx_x = Emu(int(canvas['width'] * 0.1) * 9525)
        tx_y = Emu(int(canvas['height'] * 0.1) * 9525)

    txBox = slide.shapes.add_textbox(tx_x, tx_y, tx_w, tx_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    scale = tokens.get('scale', {})
    h1_size_pt = scale.get('h1', {}).get('size', 64) // 2  # px to pt approx

    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(h1_size_pt)
    p.font.color.rgb = RGBColor.from_string(fg_hex)
    p.font.bold = True

    if subtexto:
        p2 = tf.add_paragraph()
        p2.text = ""  # spacer
        p3 = tf.add_paragraph()
        p3.text = subtexto
        p3.font.size = Pt(scale.get('body', {}).get('size', 18))
        accent_hex = palette.get('accent', fg_hex).lstrip('#')
        p3.font.color.rgb = RGBColor.from_string(accent_hex)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# ============================================================
# Fallback: SVG (sem deps externas, stdlib only)
# ============================================================

def render_svg_fallback(
    canvas_spec: dict,
    tokens: dict,
    texto: str,
    subtexto: Optional[str],
    output_path: Path,
) -> Path:
    """Gera SVG quando Pillow nao disponivel. Universal mas requer conversao para PNG/PPTX depois."""
    canvas = canvas_spec['canvas']
    palette = tokens.get('palette', {})
    bg = palette.get('neutral_light', '#FFFFFF')
    fg = palette.get('neutral_dark', '#000000')
    accent = palette.get('accent', fg)
    typo_primary = tokens.get('typography', {}).get('primary', 'serif')
    typo_secondary = tokens.get('typography', {}).get('secondary', 'sans-serif')

    safe = canvas_spec.get('ai_safe_v3') or canvas_spec.get('ai_safe_v2') or {}
    cx = canvas['width'] // 2
    cy = canvas['height'] // 2

    svg = f"""<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" width="{canvas['width']}" height="{canvas['height']}" viewBox="0 0 {canvas['width']} {canvas['height']}">
  <rect width="100%" height="100%" fill="{bg}"/>
  <text x="{cx}" y="{cy}" font-family="{typo_primary}" font-size="64" font-weight="700" fill="{fg}" text-anchor="middle" dominant-baseline="middle">{_escape_xml(texto)}</text>"""
    if subtexto:
        svg += f"""
  <text x="{cx}" y="{cy + 80}" font-family="{typo_secondary}" font-size="28" fill="{accent}" text-anchor="middle">{_escape_xml(subtexto)}</text>"""
    svg += "\n</svg>\n"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(svg, encoding="utf-8")
    return output_path


def _escape_xml(s: str) -> str:
    return (s.replace('&', '&amp;').replace('<', '&lt;')
            .replace('>', '&gt;').replace('"', '&quot;').replace("'", '&apos;'))


# ============================================================
# CLI
# ============================================================

def main() -> int:
    parser = argparse.ArgumentParser(description="Gera artes (PNG/PPTX/SVG) a partir de template + tokens + texto.")
    parser.add_argument("--template", help="Chave do template, ex: 'instagram.feed_portrait', 'linkedin.post_quadrado'")
    parser.add_argument("--texto", help="Texto principal (headline)")
    parser.add_argument("--subtexto", default=None, help="Subtexto opcional")
    parser.add_argument("--output", help="Caminho do arquivo de saida (.png/.pptx/.svg)")
    parser.add_argument("--tokens", help="Caminho para tokens.json")
    parser.add_argument("--sizes", default=None, help="Caminho para sizes.json (default: templates/arte-data/sizes.json no plugin root)")
    parser.add_argument("--layout", default="centered",
                        choices=["centered", "top-aligned", "hero-with-subtext"],
                        help="Layout do texto na AI-SAFE zone")
    parser.add_argument("--check-deps", action="store_true", help="So checa dependencias e sai")
    args = parser.parse_args()

    deps = check_deps()
    if args.check_deps:
        print_deps_status(deps)
        return 0

    # Validar args obrigatorios quando nao e --check-deps
    missing = [a for a in ("template", "texto", "output", "tokens") if not getattr(args, a)]
    if missing:
        parser.error(f"Argumentos obrigatorios faltando: {', '.join('--' + m for m in missing)}")
        return 2

    tokens_path = Path(args.tokens).resolve()
    sizes_path = Path(args.sizes).resolve() if args.sizes else (Path(__file__).parent.parent / "templates" / "arte-data" / "sizes.json")
    output_path = Path(args.output).resolve()

    tokens = load_tokens(tokens_path)
    sizes = load_sizes(sizes_path)
    canvas_spec = get_canvas_spec(sizes, args.template)

    # Decide formato pelo sufixo do output
    suffix = output_path.suffix.lower()

    try:
        if suffix == ".png":
            if not deps['pillow']:
                print("Pillow indisponivel — gerando SVG fallback (renomeie .svg)", file=sys.stderr)
                svg_path = output_path.with_suffix(".svg")
                render_svg_fallback(canvas_spec, tokens, args.texto, args.subtexto, svg_path)
                print(f"SVG fallback: {svg_path}")
            else:
                result = render_png(canvas_spec, tokens, args.texto, args.subtexto, output_path, args.layout)
                print(f"PNG gerado: {result}")
        elif suffix == ".pptx":
            if not deps['python-pptx']:
                print("python-pptx indisponivel — gerando SVG fallback (.svg)", file=sys.stderr)
                svg_path = output_path.with_suffix(".svg")
                render_svg_fallback(canvas_spec, tokens, args.texto, args.subtexto, svg_path)
                print(f"SVG fallback: {svg_path}")
            else:
                result = render_pptx(canvas_spec, tokens, args.texto, args.subtexto, output_path)
                print(f"PPTX gerado: {result}")
        elif suffix == ".svg":
            result = render_svg_fallback(canvas_spec, tokens, args.texto, args.subtexto, output_path)
            print(f"SVG gerado: {result}")
        else:
            print(f"Sufixo {suffix} nao suportado. Use .png, .pptx ou .svg", file=sys.stderr)
            return 1
    except Exception as e:
        print(f"Erro: {e}", file=sys.stderr)
        return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
